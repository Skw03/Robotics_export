#!/usr/bin/env python3

import datetime as dt
import os
import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"


def add_heading(doc, text, level=1):
    paragraph = doc.add_heading(text, level=level)
    for run in paragraph.runs:
        run.font.name = "Arial"
    return paragraph


def add_bullets(doc, items):
    for item in items:
        paragraph = doc.add_paragraph(style="List Bullet")
        paragraph.add_run(item)


def build_report():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("《机器人技术》选项A工程应用项目报告草稿")
    run.bold = True
    run.font.size = Pt(20)
    run.font.name = "Arial"
    subtitle = doc.add_paragraph("室内物流与巡检移动机器人仿真系统 | ROS 2 Humble + Gazebo Classic + Nav2")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "1. 项目概述", 1)
    doc.add_paragraph(
        "本项目面向室内物流和设施巡检场景，在 ROS 2 Humble、Gazebo Classic 与 Nav2 上实现移动机器人仿真系统。"
        "系统包含仓库和办公室两个场景，支持配送与巡检两类任务，并通过自然语言/LLM 接口映射到可执行任务。"
    )

    add_heading(doc, "2. 系统架构", 1)
    add_bullets(doc, [
        "用户命令经 CLI 或 LLM 语义层解析为 scene/task。",
        "任务通过 robotics_interfaces/Delivery action 发送到 robotics_scenario。",
        "场景管理节点加载对应行为树，将语义路线写入黑板变量。",
        "Nav2 使用 AMCL、栅格地图、代价地图、全局规划器、局部控制器和恢复行为完成导航。",
        "Gazebo 提供差速驱动、LiDAR、相机、IMU、GPS/点云等仿真传感器。"
    ])

    add_heading(doc, "3. 作业要求覆盖情况", 1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "要求"
    hdr[1].text = "当前实现"
    hdr[2].text = "后续验证"
    rows = [
        ("感知与环境建模", "LiDAR、相机、IMU、占据栅格地图、AMCL", "补 RViz/Gazebo 截图和定位稳定性观察"),
        ("运动控制与避障", "差速驱动、RPP 控制器、local costmap、collision_monitor", "baseline vs collision monitor 对比"),
        ("路径规划", "NavFn A* 与 Smac2D", "同路线多次运行，统计成功率与耗时"),
        ("应用场景", "warehouse/office；delivery/patrol", "定义成功/失败标准并记录任务日志"),
        ("人机交互", "关键词接口与 LLM 语义接口", "10 条中英文命令可靠性测试"),
        ("系统部署分析", "WSL + ROS 2 Humble 运行说明", "补硬件选型和 sim-to-real 迁移分析"),
    ]
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value

    add_heading(doc, "4. 实验计划", 1)
    add_bullets(doc, [
        "路径规划对比：NavFn A* 与 Smac2D 在同一场景、同一任务路线下运行。",
        "避障策略对比：baseline_costmap 与 collision_monitor profile 运行。",
        "LLM 可靠性：至少 10 条中英文命令，记录解析结果、延迟、置信度、失败原因。",
        "任务指标：accepted、status、task_status、elapsed_sec、route、error。"
    ])

    add_heading(doc, "5. 仿真到真机迁移分析", 1)
    doc.add_paragraph(
        "真实迁移时重点关注传感器噪声、轮式底盘打滑、执行器延迟、地图差异、动态行人、算力和网络延迟。"
        "LLM 仅作为语义规划层，不直接输出速度控制指令；安全执行仍由 Nav2、碰撞监控和人工急停保证。"
    )

    add_heading(doc, "6. AI 工具使用说明", 1)
    add_bullets(doc, [
        "Codex 用于代码生成、文档整理、实验方案制定和静态验证。",
        "OpenAI API 用于自然语言命令到 scene/task 的约束 JSON 解析。",
        "ros2-engineering-skills 用于 ROS 2/Nav2/Gazebo 工程设计检查。",
        "所有 AI 生成内容通过语法检查、dry-run、ROS action 调度和实验日志进行验证。"
    ])

    path = OUT / "robotics_option_a_final_report_draft.docx"
    doc.save(path)
    return path


def build_workbook():
    wb = Workbook()
    ws = wb.active
    ws.title = "Experiment Template"
    headers = [
        "trial", "scene", "task", "planner_profile", "avoidance_profile",
        "accepted", "status", "task_status", "elapsed_sec", "route", "error", "notes"
    ]
    ws.append(headers)
    sample_rows = [
        [1, "warehouse", "delivery", "navfn_astar", "collision_monitor", True, "SUCCEEDED", "COMPLETED", 0, "", "", "replace with real run"],
        [2, "warehouse", "delivery", "smac_2d", "collision_monitor", True, "SUCCEEDED", "COMPLETED", 0, "", "", "replace with real run"],
        [3, "office", "patrol", "navfn_astar", "baseline_costmap", False, "PENDING", "PENDING", 0, "", "", "fill after launch"],
    ]
    for row in sample_rows:
        ws.append(row)

    for col in range(1, len(headers) + 1):
        cell = ws.cell(1, col)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E78")
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        ws.column_dimensions[get_column_letter(col)].width = min(max(len(headers[col - 1]) + 4, 14), 28)
    ws.freeze_panes = "A2"

    summary = wb.create_sheet("Summary")
    summary.append(["Metric", "Formula / value"])
    summary.append(["Total trials", "=COUNTA('Experiment Template'!A2:A200)"])
    summary.append(["Completed trials", '=COUNTIF(\'Experiment Template\'!H2:H200,"COMPLETED")'])
    summary.append(["Success rate", "=IF(B2=0,0,B3/B2)"])
    summary.append(["Average elapsed sec", '=IFERROR(AVERAGEIF(\'Experiment Template\'!H2:H200,"COMPLETED",\'Experiment Template\'!I2:I200),0)'])
    for col in range(1, 3):
        summary.column_dimensions[get_column_letter(col)].width = 28
    for cell in summary[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="548235")
    summary["B4"].number_format = "0.0%"

    chart = BarChart()
    chart.title = "Course Experiment KPIs"
    chart.y_axis.title = "Value"
    chart.x_axis.title = "Metric"
    data = Reference(summary, min_col=2, min_row=2, max_row=5)
    cats = Reference(summary, min_col=1, min_row=2, max_row=5)
    chart.add_data(data, titles_from_data=False)
    chart.set_categories(cats)
    summary.add_chart(chart, "D2")

    llm = wb.create_sheet("LLM Test Cases")
    llm_headers = ["command", "expected_scene", "expected_task", "parsed_scene", "parsed_task", "parser", "latency_sec", "confidence", "correct", "failure_note"]
    llm.append(llm_headers)
    commands = [
        ("send the warehouse robot to complete a delivery loop", "warehouse", "delivery"),
        ("dispatch an office patrol route through the checkpoints", "office", "patrol"),
        ("办公室送文件到休息区", "office", "delivery"),
        ("仓库机器人巡检货架和出货区", "warehouse", "patrol"),
        ("return the office robot after checking all checkpoints", "office", "patrol"),
        ("move goods across the warehouse route", "warehouse", "delivery"),
        ("办公区补给配送", "office", "delivery"),
        ("warehouse shelf inspection loop", "warehouse", "patrol"),
        ("office mail delivery route", "office", "delivery"),
        ("仓库出库配送任务", "warehouse", "delivery"),
    ]
    for command, scene, task in commands:
        llm.append([command, scene, task, "", "", "", "", "", '=IF(AND(B2=D2,C2=E2),"yes","no")', ""])
    for row in range(2, 12):
        llm.cell(row, 9).value = f'=IF(AND(B{row}=D{row},C{row}=E{row}),"yes","no")'
    for col in range(1, len(llm_headers) + 1):
        llm.cell(1, col).font = Font(bold=True, color="FFFFFF")
        llm.cell(1, col).fill = PatternFill("solid", fgColor="7030A0")
        llm.column_dimensions[get_column_letter(col)].width = 22
    llm.column_dimensions["A"].width = 48
    llm.freeze_panes = "A2"

    path = OUT / "robotics_experiment_template.xlsx"
    wb.save(path)
    return path


def build_deck():
    extra_path = os.environ.get("PYTHON_PPTX_PATH")
    if extra_path:
        sys.path.insert(0, extra_path)
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(12, 34, 54)
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.8), Inches(1.2))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(3.0), Inches(10.8), Inches(1.0))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(210, 225, 238)
        return slide

    def content_slide(title, bullets, accent=(31, 78, 121)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(250, 252, 253)
        top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
        top.fill.solid()
        top.fill.fore_color.rgb = RGBColor(*accent)
        top.line.fill.background()
        t = slide.shapes.add_textbox(Inches(0.62), Inches(0.45), Inches(12), Inches(0.55))
        p = t.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(20, 30, 38)
        body = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(11.7), Inches(5.5))
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(35, 45, 55)
            para.space_after = Pt(10)
        return slide

    title_slide(
        "室内物流与巡检移动机器人",
        "选项A工程应用项目 | ROS 2 Humble + Gazebo Classic + Nav2"
    )
    content_slide("项目目标", [
        "在仿真环境中实现可运行的移动机器人系统，而不是只给架构图。",
        "覆盖仓库与办公室两个场景，支持配送和巡检两类任务。",
        "引入 LLM 语义层，将自然语言命令映射到可执行 ROS 2 action。"
    ])
    content_slide("系统架构", [
        "User command -> LLM/keyword parser -> Delivery.action -> robotics_scenario。",
        "行为树将语义路线转换为连续 NavigateToPose 目标。",
        "Nav2 负责 AMCL 定位、路径规划、局部控制、避障和恢复。"
    ], (84, 130, 53))
    content_slide("当前实现", [
        "两个 Gazebo 场景：AWS warehouse 与 RMF office-inspired scene。",
        "机器人模型包含差速底盘、LiDAR、相机、IMU、GPS/点云传感器。",
        "NavFn A* 与 Smac2D 已配置，collision_monitor 可用于避障对比。"
    ], (112, 48, 160))
    content_slide("任务逻辑", [
        "warehouse_delivery / warehouse_patrol / office_delivery / office_patrol 四套行为树。",
        "语义路线通过 waypoint 名称维护，便于报告和实验解释。",
        "course_experiment_runner.py 记录 accepted、status、elapsed_sec、route 和 error。"
    ])
    content_slide("LLM 语义层", [
        "course_llm_command.py 调用 OpenAI Responses API 输出受约束 JSON。",
        "输出字段固定为 scene、task、confidence、rationale。",
        "没有 API key 或调用失败时自动回退到本地关键词解析，保证演示不中断。"
    ], (192, 94, 20))
    content_slide("实验计划", [
        "路径规划对比：NavFn A* vs Smac2D。",
        "避障策略对比：baseline_costmap vs collision_monitor。",
        "LLM 可靠性：10 条中英文命令，记录正确率、延迟、失败案例。"
    ], (0, 112, 128))
    content_slide("期末交付", [
        "完成多次仿真实验并填入实验工作簿。",
        "补齐硬件选型、成本约束和 sim-to-real 迁移分析。",
        "最终报告包含 AI 工具使用说明、提示词、Skill 和验证方法。"
    ], (128, 84, 0))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(12, 34, 54)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = "下一步：运行对比实验，收集日志，完成期末答辩闭环"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    path = OUT / "robotics_option_a_midterm_deck.pptx"
    prs.save(path)
    return path


def main():
    OUT.mkdir(exist_ok=True)
    built = [build_report(), build_workbook(), build_deck()]
    manifest = OUT / "artifact_manifest.txt"
    manifest.write_text(
        "\n".join([f"{dt.datetime.now().isoformat(timespec='seconds')} {path.name}" for path in built]) + "\n",
        encoding="utf-8",
    )
    for path in built:
        print(path)


if __name__ == "__main__":
    main()
